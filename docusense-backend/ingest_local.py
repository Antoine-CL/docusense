import os, uuid, glob
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from dotenv import load_dotenv
from azure.search.documents import SearchClient
from azure.core.credentials import AzureKeyCredential
from embedding import embed_text

load_dotenv()

search = SearchClient(
    endpoint=os.getenv("AZURE_SEARCH_ENDPOINT"),
    index_name=os.getenv("AZURE_SEARCH_INDEX_NAME"),
    credential=AzureKeyCredential(os.getenv("AZURE_SEARCH_API_KEY"))
)

def extract_text(path: str) -> str:
    """Extract text from various document formats"""
    ext = os.path.splitext(path)[1].lower()
    
    try:
        if ext == ".pdf":
            pdf = fitz.open(path)
            return " ".join([page.get_text() for page in pdf])
        
        elif ext == ".docx":
            doc = Document(path)
            return "\n".join([p.text for p in doc.paragraphs])
        
        elif ext == ".pptx":
            prs = Presentation(path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            return "\n".join(text_parts)
        
        elif ext == ".txt":
            with open(path, 'r', encoding='utf-8') as f:
                return f.read()
        
        else:
            print(f"⚠️  Unsupported file type: {ext}")
            return ""
            
    except Exception as e:
        print(f"❌ Error extracting text from {path}: {e}")
        return ""

def chunk_text(text: str, size: int = 300):
    """Split text into chunks of approximately 'size' words"""
    words = text.split()
    for i in range(0, len(words), size):
        chunk_text = " ".join(words[i:i+size])
        if chunk_text.strip():  # Only yield non-empty chunks
            yield chunk_text, i // size

def ingest_documents(folder_path: str = "sample_docs"):
    """Ingest all documents from the specified folder"""
    if not os.path.exists(folder_path):
        print(f"📁 Creating {folder_path} folder...")
        os.makedirs(folder_path)
        print(f"📝 Please add some documents to {folder_path}/ and run again")
        return
    
    docs = []
    file_count = 0
    chunk_count = 0
    
    # Process all supported files in the folder
    for file_path in glob.glob(f"{folder_path}/*"):
        if os.path.isfile(file_path):
            print(f"📄 Processing: {os.path.basename(file_path)}")
            
            content = extract_text(file_path)
            if not content.strip():
                print(f"⚠️  No text extracted from {file_path}")
                continue
                
            title = os.path.basename(file_path)
            file_count += 1
            
            # Split into chunks and create embeddings
            for snippet, chunk_idx in chunk_text(content):
                try:
                    embedding = embed_text(snippet)
                    
                    doc = {
                        "id": str(uuid.uuid4()),
                        "title": title,
                        "content": snippet,
                        "chunk": chunk_idx,
                        "vector": embedding
                    }
                    
                    docs.append(doc)
                    chunk_count += 1
                    
                    # Upload in batches of 100 to avoid timeouts
                    if len(docs) >= 100:
                        search.upload_documents(documents=docs)
                        print(f"📤 Uploaded batch of {len(docs)} chunks")
                        docs = []
                        
                except Exception as e:
                    print(f"❌ Error processing chunk from {title}: {e}")
    
    # Upload remaining documents
    if docs:
        search.upload_documents(documents=docs)
        print(f"📤 Uploaded final batch of {len(docs)} chunks")
    
    print(f"✅ Ingestion complete!")
    print(f"📊 Processed {file_count} files, created {chunk_count} chunks")

if __name__ == "__main__":
    ingest_documents() 